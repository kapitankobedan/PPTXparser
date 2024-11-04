import os
import json
from pptx import Presentation
import fitz  # PyMuPDF

PATH = r'C:\Users\Данила\Downloads\1_lecture.pdf'

# Функция для создания папки, если её не существует
def create_output_directory(output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)


# Парсинг PPTX
def parse_pptx(pptx_file, output_dir):
    prs = Presentation(pptx_file)

    presentation_data = {
        "title": prs.core_properties.title or os.path.basename(pptx_file),  # Тема презентации
        "slides": []
    }

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "title": "",
            "text": "",
            "images": []
        }

        for shape in slide.shapes:
            # Извлечение текста
            if hasattr(shape, "text"):
                if shape.text_frame is not None and shape.text_frame.paragraphs:
                    if not slide_data["title"]:  # Первый текст на слайде как заголовок
                        slide_data["title"] = shape.text.strip()
                    else:
                        slide_data["text"] += shape.text.strip() + "\n"

            # Извлечение изображений
            if shape.shape_type == 13:  # Если тип объекта - изображение
                image = shape.image
                image_filename = f"slide_{slide_num}_image_{shape.shape_id}.jpg"
                image_path = os.path.join(output_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image.blob)
                slide_data["images"].append(image_filename)

        presentation_data["slides"].append(slide_data)

    return presentation_data


# Парсинг PDF
def parse_pdf(pdf_file, output_dir):
    pdf_document = fitz.open(pdf_file)

    presentation_data = {
        "title": os.path.basename(pdf_file),  # Тема презентации — имя файла
        "slides": []
    }

    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)

        # Получаем текст страницы в виде блоков с информацией о шрифтах
        blocks = page.get_text("dict")["blocks"]
        slide_data = {
            "title": "",
            "text": "",
            "images": []
        }

        # Попытка выделить первую строку как заголовок
        if blocks:
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        for span in line["spans"]:
                            if not slide_data["title"]:
                                # Первая строка считается заголовком
                                slide_data["title"] = span["text"].strip()
                            else:
                                # Все последующие строки считаются основным текстом
                                slide_data["text"] += span["text"].strip() + " "
                        slide_data["text"] += "\n"

        # Извлечение изображений, только для текущей страницы
        images = page.get_images(full=True)
        for image_index, img in enumerate(images):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_filename = f"slide_{page_num + 1}_image_{image_index}.{image_ext}"
            image_path = os.path.join(output_dir, image_filename)

            with open(image_path, "wb") as image_file:
                image_file.write(image_bytes)

            slide_data["images"].append(image_filename)

        # Добавляем слайд с данными о заголовке, тексте и изображениях
        presentation_data["slides"].append(slide_data)

    pdf_document.close()

    return presentation_data


# Основная функция для парсинга презентации
def parse_presentation(file_path):
    output_dir = f"{os.path.splitext(os.path.basename(file_path))[0]}_output"  # Создаем имя папки на основе имени файла
    create_output_directory(output_dir)

    extension = os.path.splitext(file_path)[1].lower()

    if extension == ".pptx":
        presentation_data = parse_pptx(file_path, output_dir)
    elif extension == ".pdf":
        presentation_data = parse_pdf(file_path, output_dir)
    else:
        raise ValueError("Unsupported file format")

    # Сохранение данных в JSON
    json_file = os.path.join(output_dir, "presentation_data.json")
    save_to_json(presentation_data, json_file)

    print(f"Данные сохранены в папку: {output_dir}")
    return presentation_data


# Функция для сохранения данных в JSON-файл
def save_to_json(data, output_file):
    with open(output_file, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=4)


if __name__ == "__main__":
    file_path = PATH
    parse_presentation(file_path)

